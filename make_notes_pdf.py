"""Compose the portrait notes-pages PDF (slide on top, speaker notes below).

Requires the slide PNGs exported by PowerPoint COM into presentation/slides/
(see repo notes) and the pptx built by build_presentation.py.
Usage:  python -u make_notes_pdf.py
"""
import os, textwrap
from pptx import Presentation
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.image as mpimg
from matplotlib.backends.backend_pdf import PdfPages

REPO = os.path.dirname(os.path.abspath(__file__))
PRES = os.path.join(REPO, 'presentation')
prs = Presentation(os.path.join(PRES, 'finetuning_results.pptx'))
notes = [s.notes_slide.notes_text_frame.text if s.has_notes_slide else ''
         for s in prs.slides]
out = os.path.join(PRES, 'finetuning_results_notes.pdf')
with PdfPages(out) as pdf:
    for i, note in enumerate(notes, 1):
        fig = plt.figure(figsize=(8.27, 11.69))
        ax = fig.add_axes([0.06, 0.55, 0.88, 0.40])
        ax.imshow(mpimg.imread(os.path.join(PRES, 'slides',
                                            f'slide_{i:02d}.png')))
        ax.axis('off')
        ax.set_anchor('N')
        axn = fig.add_axes([0.08, 0.04, 0.84, 0.48])
        axn.axis('off')
        wrapped = '\n'.join(textwrap.fill(p, 95) for p in note.split('\n'))
        axn.text(0, 1, wrapped, va='top', ha='left', fontsize=10.5,
                 family='sans-serif', linespacing=1.5)
        fig.text(0.5, 0.015, f'{i} / {len(notes)}', ha='center', fontsize=9,
                 color='#898781')
        pdf.savefig(fig)
        plt.close(fig)
print(out)
