"""
Build the project presentation (for David & Shane): hypothesis, tests,
results, implications — with full speaker notes on every slide.

Output: presentation/finetuning_results.pptx
Usage:  python -u build_presentation.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

REPO = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(REPO, 'figures')
OUT = os.path.join(REPO, 'presentation')
os.makedirs(OUT, exist_ok=True)

INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
BLUE = RGBColor(0x2A, 0x78, 0xD6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(title, bullets=None, image=None, notes='', img_width=11.6,
              title_size=30):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3),
                              Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = INK
    top = 1.25
    if bullets:
        bb = s.shapes.add_textbox(Inches(0.7), Inches(top), Inches(12.0),
                                  Inches(5.6))
        tf = bb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = b
            para.font.size = Pt(20)
            para.font.color.rgb = INK2 if b.startswith(' ') else INK
            para.space_after = Pt(12)
        top = None
    if image:
        path = os.path.join(FIG, image)
        from PIL import Image
        w, h = Image.open(path).size
        disp_w = img_width
        disp_h = disp_w * h / w
        if disp_h > 5.9:
            disp_h = 5.9
            disp_w = disp_h * w / h
        left = (13.333 - disp_w) / 2
        s.shapes.add_picture(path, Inches(left), Inches(1.3),
                             width=Inches(disp_w))
    s.notes_slide.notes_text_frame.text = notes
    return s


# ---------------------------------------------------------------- 1 title
add_slide(
    'Choose your teachers:\nwhat happened when we tried to fine-tune the '
    'eruption forecaster',
    bullets=['Marsden fine-tuning project — results update',
             'Alberto Ardid · September 2026',
             ' Phases A–D + G complete · all results pseudo-prospectively '
             'tested · code + results on GitHub (aardid/puia_fine_tunning_dev)'],
    title_size=32,
    notes=(
        "Thanks for making time. Over the last few days I pushed the Marsden "
        "fine-tuning project through essentially the whole planned arc — the "
        "leakage-audited baseline, source-pool curation, SER and STRUT tree "
        "refinement, TransferBoost, and the pseudo-prospective unified "
        "evaluation. The punchline is not what we proposed: fine-tuning the "
        "model turns out to be the wrong lever at data-scarce volcanoes. "
        "Choosing WHICH volcanoes we train on matters far more — and one "
        "volcano, Whakaari, is actively poisoning the pool. I want to walk "
        "you through the hypothesis, the tests, the results, and what I "
        "think it means for the paper and for observatories. Everything "
        "here is scripted and reproducible; every number I show survived a "
        "train-on-the-past-only replay unless I flag otherwise."))

# ---------------------------------------------------------------- 2 question
add_slide(
    'The question we started from',
    bullets=[
        '2025 paper: precursors are ergodic — a pool model forecasts unseen '
        'volcanoes (AUC ≈ 0.8)',
        'Marsden question: can we FINE-TUNE the generalized model to a '
        'specific volcano?',
        'Hypothesis: adapting the trees to local data recovers the gap to '
        'direct training',
        ' Targets: Ruapehu (3 eruptions), Tongariro (2) — later Whakaari (5)',
        ' Pool: Whakaari, Ruapehu, Tongariro, Ontake, St Helens '
        '(12 phreatic eruptions)'],
    notes=(
        "Quick framing. Our 2025 paper showed eruption precursors are "
        "ergodic: train on many volcanoes, forecast one you've never seen, "
        "and you get roughly the same skill as training directly on it. The "
        "Marsden proposal asked the natural next question: if we have SOME "
        "local data at a target volcano, can we fine-tune the generalized "
        "model to do better there? Our working hypothesis was yes — that "
        "adapting the ensemble's trees with target data would add skill, "
        "with Ruapehu and Tongariro as the data-scarce targets. The plan "
        "was phased: an honest baseline first, then pool curation as a "
        "cheap alternative, then the fine-tuning methods proper — SER and "
        "STRUT tree refinement, and TransferBoost. Spoiler: the hypothesis "
        "was wrong in an interesting way, and the cheap alternative turned "
        "out to be the headline."))

# ---------------------------------------------------------------- 3 baseline
add_slide(
    'First: an honest baseline (Phase A)',
    bullets=[
        'Fixed a critical label-alignment bug (silent NaN→True: 78% of '
        'labels were wrong)',
        'Fixed the crashed pipeline + filled missing forecast years',
        'Leave-one-eruption-out over the pool: AUC = 0.823',
        ' Matches the 2025 paper (~0.8) — replication gate passed',
        ' Everything below runs locally from cached features — no server '
        'needed'],
    notes=(
        "Before testing anything I made the baseline trustworthy. Two things "
        "mattered. First, there was a nasty silent bug in train_one_model: "
        "after undersampling, a pandas reindex turned NaNs into TRUE labels, "
        "so about three quarters of training labels were positive — the "
        "model was degenerate. Fixed. Second, the Phase A run had been stuck "
        "since May on a path bug, and the model-00 forecasts were missing "
        "years without anyone noticing. With all that repaired, the "
        "leave-one-eruption-out baseline over the five-volcano phreatic pool "
        "comes out at 0.823 — consistent with the paper. That's our "
        "replication gate, and it passed. One nice practical discovery: "
        "because all the tsfresh features are cached, every experiment I'll "
        "show ran on my desktop against the network drive — training "
        "included. No server queue."))

# ---------------------------------------------------------------- 4 test1
add_slide(
    'Test 1 — is the training pool itself the problem? (Phase B)',
    bullets=[
        'Before touching the model: which source volcanoes help, which harm?',
        'Trained ALL 15 source-pool subsets per target '
        '(leave-target-volcano-out)',
        'Scored on the target\'s full record — target never seen in training',
        ' Also: nested selection (pool chosen without peeking at the test '
        'eruption)'],
    notes=(
        "The first test doesn't touch the model at all. The transfer-learning "
        "literature says dropping dissimilar source tasks often beats "
        "fine-tuning on the target — so before adapting anything, I asked "
        "which of our source volcanoes actually help. The test is brute "
        "force: for each target, train the full ensemble on every one of the "
        "fifteen possible subsets of the other four volcanoes, then forecast "
        "the target's entire record. The target never appears in training, "
        "so every number is fully out-of-sample. And to keep ourselves "
        "honest, I also ran nested selection — for each eruption, choosing "
        "the best pool using only the OTHER eruptions — so the pool choice "
        "itself is never informed by the event it's scored on."))

# ---------------------------------------------------------------- 5 result1
add_slide('Result: every Whakaari-free pool beats every '
          'Whakaari-containing pool',
          image='F1_pool_landscape.png',
          notes=(
              "Here's the landscape, and it stunned me. Blue pools exclude "
              "Whakaari, orange pools include it. At all three targets — "
              "Ruapehu, Tongariro, and Whakaari itself as a target — every "
              "single blue pool outperforms every single orange pool. "
              "Perfect rank separation, forty-five comparisons, no overlap. "
              "Ruapehu's best forecaster is St Helens alone at 0.96. "
              "Tongariro's is Ontake plus St Helens at 0.89 — versus 0.65 "
              "for the full pool. And the same tiny foreign pool tops the "
              "table for Whakaari at 0.97. Note the irony: Whakaari "
              "contributes five of our twelve training eruptions — the "
              "most data — and it is toxic everywhere. The nested, "
              "no-peeking version holds up: 0.95, 0.88, 0.86 — the "
              "selection is stable and the bias penalty is under 0.02. "
              "This is not a fluke of hindsight."))

# ---------------------------------------------------------------- 6 external
add_slide('It generalizes: 8 of 8 volcanoes the models never saw',
          image='F2_external_validation.png',
          notes=(
              "Obvious objection: maybe that's a New Zealand quirk. So I "
              "took the trained ensembles, unchanged, to eight volcanoes "
              "completely outside the pool — Pavlof twice, Veniaminof, "
              "Bezymianny, Copahue, Montserrat, and two more — thirty-one "
              "eruptions no model had ever seen. The curated "
              "Tongariro-Ontake-St Helens pool beats the full five-source "
              "generalized model on eight out of eight. Sign test p is "
              "about 0.004. And in the matched pairs where the only "
              "difference is Whakaari's presence, dropping it helps at all "
              "eight. Copahue forecasts at 0.95, Bezymianny at 0.93, from "
              "purely foreign data. So the curation rule we learned in New "
              "Zealand transfers across three continents and multiple "
              "eruption styles. This, for me, is what makes the result "
              "paper-grade rather than a local observation."))

# ---------------------------------------------------------------- 7 mechanism
add_slide('Why: loudness does not transfer — waveform shape does',
          image='F3_feature_mechanism.png',
          notes=(
              "So why would the most eruption-rich volcano poison training? "
              "I tallied what features the trees actually select. With "
              "Whakaari in the pool, selection locks onto amplitude: "
              "seventy-one percent RSAM features, FFT power, amplitude "
              "quantiles — because Whakaari's twenty positive windows are "
              "loud and they dominate the Mann-Whitney ranking. Remove "
              "Whakaari, and selection shifts to temporal structure — "
              "autocorrelation, wavelet coefficients — spread across the "
              "mid and high frequency bands. And that's the physics of the "
              "whole story: every volcano's 'loud' is different, so "
              "absolute amplitude doesn't transfer; the SHAPE of the "
              "pre-eruptive ramp is what's shared — that's the ergodic "
              "part. It also explains the asymmetry we see: Whakaari's "
              "data harms others, yet Whakaari itself is the EASIEST "
              "target to forecast from foreign, shape-based models. Shane — "
              "I'd value your read on the volcanological interpretation "
              "here."))

# ---------------------------------------------------------------- 8 finetune
add_slide('Test 2 — actual fine-tuning (SER/STRUT tree refinement, Phase C)',
          image='F5_decision_ladder.png',
          notes=(
              "Now the thing we actually proposed: fine-tuning. I "
              "implemented SER and STRUT — Segev's tree-refinement transfer "
              "methods — which adapt the ensemble's split thresholds and "
              "structure using target data, evaluated leave-one-eruption-"
              "out. The result is this ladder. At Tongariro and Ruapehu — "
              "two and three eruptions — refinement is unreliable and "
              "mostly destructive: the trees memorize the few training "
              "eruptions and go silent on the held-out one; the worst case "
              "loses a third of the AUC. At Whakaari, with five eruptions, "
              "every method helps — up to plus 0.08. Two technical points: "
              "per-tree undersampling of the refinement data was essential "
              "— without it the ensemble collapses to clones — and gentle "
              "variants like calibration-only don't pay at all. So there's "
              "a data threshold somewhere between three and five recorded "
              "eruptions below which fine-tuning actively hurts. That's "
              "our hypothesis, falsified in the regime we proposed it "
              "for — and I think that negative result is publishable on "
              "its own."))

# ---------------------------------------------------------------- 9 tboost
add_slide(
    'Test 3 — automatic instance weighting (TrAdaBoost, Phase D)',
    bullets=[
        'Hope: boosting automatically downweights unhelpful source data',
        'Reality 1: weight collapses onto the target (no source transfer)',
        'Reality 2: it cannot see the Whakaari harm — WIZ keeps ~50% of '
        'residual source weight',
        'Reality 3: retrospective AUC ≈ 0.997 at all targets… ',
        ' …which the prospective test exposed as sibling-eruption '
        'memorization (next slide)'],
    notes=(
        "Third test: TrAdaBoost — instance-weighted boosting that's "
        "supposed to automatically downweight unhelpful source data, the "
        "automatic version of our pool curation. Three findings. One, the "
        "classic pathology: training weight collapses almost entirely onto "
        "the target within forty iterations, so no real transfer happens. "
        "Two — and this is the interesting one — it is completely blind to "
        "the Whakaari problem. Whakaari keeps about half the remaining "
        "source weight, because its windows are individually EASY to "
        "classify. The harm operates through feature selection, invisible "
        "at the level of single data points. You cannot fix this dataset "
        "problem with instance weighting. Three, the trap: retrospectively "
        "it scored 0.997 everywhere — looks like a miracle. I didn't "
        "believe it, and the next slide shows what the honest test did to "
        "that number."))

# ---------------------------------------------------------------- 10 phase G
add_slide('The referee: everything replayed with only prior information '
          '(Phase G)',
          image='F4_prospective_scoreboard.png',
          notes=(
              "This is the slide that disciplines all the others. Every "
              "method, rebuilt for every eruption using only information "
              "available a month before it — data, pool choice, "
              "adaptation, everything — and scored on the hard unrest-month "
              "background. Three verdicts. Curation survives: the "
              "prospectively-selected pool beats the full pool at all "
              "three targets, even when early history gives it only one or "
              "two eruptions to choose from. Tree refinement does NOT "
              "survive at two-to-three eruptions — consistent with the "
              "ladder. And tboost's 0.997 collapses to literal zero at "
              "every genuinely prospective Ruapehu and Tongariro eruption "
              "— it had been recognizing sibling eruptions from the same "
              "unrest episodes. But — at data-rich Whakaari it posts 0.91, "
              "the best prospective score in the study, including "
              "detecting December 2019 at full confidence from four prior "
              "eruptions. The asterisked bars are fallback-carried; the "
              "honest per-eruption replay is in backup."))

# ---------------------------------------------------------------- 11 dec2019
add_slide('What this looks like in operations: Whakaari, December 2019',
          image='M4_operational_alerts.png',
          notes=(
              "And here's why anyone outside our field should care. Top "
              "panel: two years of the Ontake-trained model at Whakaari "
              "under a causal alert rule I then optimized — threshold is "
              "the trailing 180-day 99.5th percentile, sustained a full "
              "day. Three false alarms in two restless years, and the "
              "final alert begins nine days before the December 9 "
              "eruption. Bottom left, the final month: the sustained "
              "escalation, plus the fully prospective target-boosted "
              "model firing at 3.4 days. Bottom right, the honest cost "
              "accounting: the full detection versus false-alarm frontier "
              "across 240 rule settings — one eruption in five at one "
              "false alarm a year, three in five at six. Important "
              "honesty: the rule tuning is in-sample, and the nine-day "
              "lead is robust across most settings, but before we say "
              "this publicly I want the parameters fixed on pre-2019 data "
              "and 2019 verified untouched. That run is queued. Still — "
              "for GNS and the Marsden panel, this is the picture."))

# ---------------------------------------------------------------- 12 ladder
add_slide(
    'Implications: an evidence-based recipe by eruption history',
    bullets=[
        '0 recorded eruptions → curated foreign pool + calibrate threshold '
        'to local noise (free +0.03 pooled)',
        '1–3 eruptions → use them ONLY to select source volcanoes '
        '(+0.10 to +0.24); do not retrain',
        '~5+ eruptions → target-boosted local model finally wins '
        '(0.91 prospective at Whakaari)',
        ' And one rule everywhere: leave Whakaari out of phreatic '
        'training pools'],
    notes=(
        "Pulling it together, the practical output is a decision ladder an "
        "observatory can actually follow, and every rung was verified "
        "under the prospective replay. If a volcano has no recorded "
        "eruptions: run the curated foreign-pool model, and calibrate the "
        "alert threshold to the local noise floor — that calibration alone "
        "recovered three points of pooled AUC and needs only quiet data. "
        "With one to three eruptions: spend them selecting your source "
        "volcanoes — worth ten to twenty-four points — and resist "
        "retraining, which we showed is actively destructive there. From "
        "about five eruptions: train on local data; boosting finally "
        "outruns transfer. And across the board: Whakaari's data stays "
        "out of the pool. Its amplitude-dominated signature is a "
        "cautionary tale about 'more data is better' — the single most "
        "counterintuitive and, I think, most citable sentence of the "
        "study."))

# ---------------------------------------------------------------- 13 caveats
add_slide(
    'Caveats, and what remains before submission',
    bullets=[
        'Eruption counts are small (2–5 per target) — TPR is coarse; '
        'aggregate stats carry the claims',
        'Alert-rule tuning is in-sample → pre-2019 validation run queued',
        'Bootstrap CIs on headline AUCs still to add (PLAN commitment)',
        'Style-dependence: VRLE/Veniaminof forecast poorly — Phase H to map '
        'where transfer works',
        ' Proposed venue: Nature Communications — '
        '"Choosing training volcanoes beats fine-tuning…"'],
    notes=(
        "The honest limits. First, per-target eruption counts are tiny, so "
        "individual AUCs move in coarse steps — the claims rest on the "
        "aggregate patterns: the perfect 45-comparison separation, the "
        "8-of-8 external wins, the prospective replay. Second, the alert "
        "rule in the 2019 case study is tuned in-sample; the pre-2019 "
        "fix-and-verify run is the next thing I'll do. Third, PLAN.md "
        "promised bootstrap confidence intervals on the headline numbers — "
        "straightforward to add. Fourth, transfer isn't universal: "
        "Veniaminof and VRLE forecast poorly from any pool, and mapping "
        "that style-dependence is the remaining Phase H work — it likely "
        "sharpens the phreatic-versus-magmatic boundary of the ergodic "
        "claim. Given all that, my proposal is Nature Communications as "
        "the natural sequel to the 2025 paper, framed around curation "
        "beating fine-tuning. I'd love your take on framing, authorship, "
        "and whether GNS should see the December 2019 material before "
        "submission."))

# ---------------------------------------------------------------- 14 discuss
add_slide(
    'Discussion',
    bullets=[
        'Is the mechanism story (amplitude vs. shape) volcanologically '
        'sound? (Shane)',
        'Prospective protocol tight enough for review? Anything I missed? '
        '(David)',
        'Do we split the operational/alert material or keep one paper?',
        'Whakaari sensitivities: how do we frame Dec 2019 responsibly?',
        ' Everything reproducible: one repo, one figure script, all phases '
        'scripted'],
    notes=(
        "Some things I genuinely want your input on. Shane — the mechanism: "
        "does amplitude-versus-waveform-shape hold up against what we know "
        "about Whakaari's shallow hydrothermal system and Ontake and St "
        "Helens as analogues? Is there a physical reason those two carry "
        "such clean precursor structure? David — the evaluation: is the "
        "pseudo-prospective protocol airtight enough for referees, and do "
        "you see any leakage channel I've missed — window overlap at "
        "cutoffs, the source pools containing post-date data, anything? "
        "Then two strategic calls: one paper or two — the science paper "
        "plus a Bulletin of Volcanology operational piece — and how we "
        "frame the December 2019 case study responsibly, given the "
        "sensitivities and the inquest history. Everything you've seen "
        "regenerates from one repository — happy to walk through any "
        "number live."))

path = os.path.join(OUT, 'finetuning_results.pptx')
prs.save(path)
print(path)
